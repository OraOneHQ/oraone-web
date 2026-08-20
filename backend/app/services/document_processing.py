"""Document processing (Phase 7).

Two layers:

1. **Pure functions** — extract text per filetype + chunk it with overlap.
   Easy to unit-test, no IO unless you give them bytes.
2. **Pipeline orchestrator** — loads bytes from storage, runs the parser
   + chunker, persists chunks, flips ``documents.status`` and writes
   processing telemetry.

Supported file types:

* PDF  (``application/pdf``, ``*.pdf``)
* DOCX (``*.docx``)
* TXT  (``text/plain``, ``*.txt``)
* MD   (``*.md``, ``text/markdown``)
* CSV  (``*.csv``)

Phase 9: after chunking, each chunk is embedded (Amazon Titan v2 or the
hashing fallback) and the vector is stored in ``document_chunks.embedding``
so the RAG retriever can do cosine-similarity search. The
``DocumentChunk.metadata`` JSONB stores citation hints (page number,
section heading, source file).
"""
from __future__ import annotations

import csv
import io
import logging
import re
import uuid
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Iterable, Optional

from sqlalchemy import delete, select
from sqlalchemy.ext.asyncio import AsyncSession

from app.database.models.document import Document, DocumentStatus
from app.database.models.document_chunk import DocumentChunk
from app.database.session import AsyncSessionLocal, init_engine
from app.providers.embeddings import get_embedding_provider
from app.services import storage

log = logging.getLogger("app.doc_processing")

# Chunking knobs from the spec
TARGET_CHUNK_CHARS = 800
CHUNK_OVERLAP_CHARS = 100


# ────────────────────────── text extraction ──────────────────────────

@dataclass
class ExtractedPage:
    """One logical "page" of source text. PDFs have real pages; for
    flat formats we synthesise a single page so the chunker sees a
    uniform input shape."""

    page: int
    text: str
    section: Optional[str] = None


def extract_text(filename: str, content_type: Optional[str], data: bytes) -> list[ExtractedPage]:
    """Dispatch to the right parser based on filename / content type.

    Returns one or more ``ExtractedPage`` blocks. Raises
    ``ValueError`` if the file type isn't supported.
    """
    ext = Path(filename).suffix.lower().lstrip(".")
    ct = (content_type or "").lower()

    if ext == "pdf" or ct == "application/pdf":
        return _extract_pdf(data)
    if ext == "docx" or ct in (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ):
        return _extract_docx(data)
    if ext in ("xlsx", "xlsm") or ct in (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ):
        return _extract_xlsx(data)
    if ext == "pptx" or ct in (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ):
        return _extract_pptx(data)
    if ext == "json" or ct == "application/json":
        return _extract_json(data)
    if ext in ("html", "htm") or ct == "text/html":
        return _extract_html(data)
    if ext == "csv" or ct == "text/csv":
        return _extract_csv(data)
    if ext == "md" or ct == "text/markdown":
        return _extract_text(data, ext="md")
    if ext == "txt" or ct.startswith("text/"):
        return _extract_text(data, ext="txt")

    raise ValueError(
        f"Unsupported file type for chunking: filename={filename!r} content_type={content_type!r}"
    )


def _decode(data: bytes) -> str:
    # Try UTF-8 first, then latin-1 as a forgiving fallback — never raise.
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="replace")


def _extract_text(data: bytes, *, ext: str) -> list[ExtractedPage]:
    text = _decode(data)
    section = None
    if ext == "md":
        # Pick the first H1/H2 heading as the section label, if any.
        m = re.search(r"^#{1,2}\s+(.+?)\s*$", text, flags=re.MULTILINE)
        if m:
            section = m.group(1).strip()
    return [ExtractedPage(page=1, text=text, section=section)]


def _extract_csv(data: bytes) -> list[ExtractedPage]:
    """Flatten a CSV into a single text block: ``col1: val1 | col2: val2`` per row.

    This is dumb on purpose — the goal is "the file is now in the chunker"
    not "perfect tabular retrieval". A later phase can replace this with a
    proper table-aware embedding.
    """
    text = _decode(data)
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return [ExtractedPage(page=1, text="")]

    headers = rows[0]
    out_lines: list[str] = []
    for row in rows[1:]:
        parts = []
        for i, cell in enumerate(row):
            header = headers[i] if i < len(headers) else f"col{i + 1}"
            parts.append(f"{header}: {cell}")
        out_lines.append(" | ".join(parts))

    return [ExtractedPage(page=1, text="\n".join(out_lines), section="CSV rows")]


def _extract_pdf(data: bytes) -> list[ExtractedPage]:
    """Per-page PDF text extraction.

    ``pypdf`` is pure-Python and works on the kinds of PDFs we expect
    (digital, not scanned). Scanned PDFs would need OCR — out of scope
    for Phase 7 by design.
    """
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages: list[ExtractedPage] = []
    for i, p in enumerate(reader.pages, start=1):
        try:
            text = p.extract_text() or ""
        except Exception as e:  # malformed page — skip but keep going
            log.warning("pdf_extract page=%d failed: %s", i, e)
            text = ""
        pages.append(ExtractedPage(page=i, text=text))
    return pages


def _extract_docx(data: bytes) -> list[ExtractedPage]:
    """Paragraph-level DOCX extraction. We keep the running ``section``
    pointer at the most recent heading paragraph so each chunk inherits
    its enclosing section.
    """
    from docx import Document as DocxDocument  # python-docx

    doc = DocxDocument(io.BytesIO(data))
    blocks: list[str] = []
    sections: list[Optional[str]] = []
    current_section: Optional[str] = None

    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name or "").lower() if para.style else ""
        if "heading" in style or "title" in style:
            current_section = text
        blocks.append(text)
        sections.append(current_section)

    if not blocks:
        return [ExtractedPage(page=1, text="")]

    # DOCX has no real page concept here — emit one ExtractedPage with
    # all paragraphs joined and the *first* non-None section as the label.
    joined = "\n".join(blocks)
    label = next((s for s in sections if s), None)
    return [ExtractedPage(page=1, text=joined, section=label)]


def _extract_xlsx(data: bytes) -> list[ExtractedPage]:
    """One ExtractedPage per worksheet. Rows are flattened to
    ``cell | cell | cell`` lines so the chunker sees plain text."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    pages: list[ExtractedPage] = []
    for i, ws in enumerate(wb.worksheets, start=1):
        lines: list[str] = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
        if lines:
            pages.append(ExtractedPage(page=i, text="\n".join(lines), section=ws.title))
    wb.close()
    return pages or [ExtractedPage(page=1, text="")]


def _extract_pptx(data: bytes) -> list[ExtractedPage]:
    """One ExtractedPage per slide; concatenates all shape text frames."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    pages: list[ExtractedPage] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        title = None
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text)
            if text.strip():
                parts.append(text)
                if title is None:
                    title = text.splitlines()[0][:120]
        if parts:
            pages.append(ExtractedPage(page=i, text="\n".join(parts), section=title))
    return pages or [ExtractedPage(page=1, text="")]


def _extract_json(data: bytes) -> list[ExtractedPage]:
    """Flatten JSON into ``path: value`` lines so nested config / API
    payloads become searchable text. Falls back to raw text on parse error."""
    import json as _json

    raw = _decode(data)
    try:
        obj = _json.loads(raw)
    except Exception:
        return [ExtractedPage(page=1, text=raw)]

    lines: list[str] = []

    def walk(node, prefix: str = "") -> None:
        if isinstance(node, dict):
            for k, v in node.items():
                walk(v, f"{prefix}.{k}" if prefix else str(k))
        elif isinstance(node, list):
            for idx, v in enumerate(node):
                walk(v, f"{prefix}[{idx}]")
        else:
            lines.append(f"{prefix}: {node}")

    walk(obj)
    return [ExtractedPage(page=1, text="\n".join(lines), section="JSON")]


def _extract_html(data: bytes) -> list[ExtractedPage]:
    """Strip tags to plain text using the stdlib HTML parser (no external
    dependency). Script/style content is dropped; the <title> becomes the
    section label."""
    from html.parser import HTMLParser

    class _Stripper(HTMLParser):
        def __init__(self) -> None:
            super().__init__()
            self.parts: list[str] = []
            self.title: Optional[str] = None
            self._skip = 0
            self._in_title = False

        def handle_starttag(self, tag, attrs):
            if tag in ("script", "style"):
                self._skip += 1
            elif tag == "title":
                self._in_title = True
            elif tag in ("p", "br", "div", "li", "tr", "h1", "h2", "h3", "h4"):
                self.parts.append("\n")

        def handle_endtag(self, tag):
            if tag in ("script", "style") and self._skip:
                self._skip -= 1
            elif tag == "title":
                self._in_title = False

        def handle_data(self, data):
            if self._skip:
                return
            text = data.strip()
            if not text:
                return
            if self._in_title and not self.title:
                self.title = text[:160]
            self.parts.append(text + " ")

    parser = _Stripper()
    parser.feed(_decode(data))
    return [ExtractedPage(page=1, text="".join(parser.parts), section=parser.title)]


# ────────────────────────── normalization ──────────────────────────

_WS_RE = re.compile(r"[ \t\u00A0]+")  # spaces, tabs, NBSP
_NL_RE = re.compile(r"\n{3,}")       # collapse 3+ newlines to 2


def normalize(text: str) -> str:
    """Whitespace normalisation that preserves paragraph breaks."""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = _WS_RE.sub(" ", text)
    text = _NL_RE.sub("\n\n", text)
    return text.strip()


# ────────────────────────── chunking ──────────────────────────

# Heading regexes used to bias chunk boundaries (purely heuristic).
_MD_HEADING_RE = re.compile(r"^#{1,6}\s+(.+?)\s*$", re.MULTILINE)
_NUMBERED_HEADING_RE = re.compile(r"^\s*\d+(\.\d+)*\s+[A-Z][^\n]{2,80}$", re.MULTILINE)


@dataclass
class Chunk:
    """In-memory representation of one chunk before it's persisted."""

    index: int
    content: str
    metadata: dict


def _detect_section(text: str, fallback: Optional[str]) -> Optional[str]:
    """Best-effort: pick the most recent heading-looking line above ``text``."""
    for rx in (_MD_HEADING_RE, _NUMBERED_HEADING_RE):
        matches = list(rx.finditer(text))
        if matches:
            return matches[-1].group(0).lstrip("#").strip()
    return fallback


def chunk_pages(
    pages: Iterable[ExtractedPage],
    *,
    source_file: str,
    target_chars: int = TARGET_CHUNK_CHARS,
    overlap_chars: int = CHUNK_OVERLAP_CHARS,
) -> list[Chunk]:
    """Slide a window of ``target_chars`` with ``overlap_chars`` carry-over.

    Returns a flat list of ``Chunk`` objects ordered by appearance, with
    each chunk's metadata stamped with ``page`` / ``section`` /
    ``source_file``. Empty inputs return an empty list (callers can
    decide whether that's an error).
    """
    if target_chars <= 0:
        raise ValueError("target_chars must be positive")
    if overlap_chars < 0 or overlap_chars >= target_chars:
        raise ValueError("overlap_chars must satisfy 0 <= overlap < target_chars")

    chunks: list[Chunk] = []
    idx = 0
    stride = target_chars - overlap_chars

    for page in pages:
        text = normalize(page.text or "")
        if not text:
            continue

        # Walk the page text with the sliding window.
        i = 0
        page_section_fallback = page.section
        while i < len(text):
            end = min(i + target_chars, len(text))
            # Prefer to end on a paragraph break or sentence boundary
            # within the last 20% of the window — purely a quality hint.
            if end < len(text):
                tail = text[i:end]
                snap = max(
                    tail.rfind("\n\n"),
                    tail.rfind(". "),
                    tail.rfind("! "),
                    tail.rfind("? "),
                )
                soft_floor = int(target_chars * 0.8)
                if snap >= soft_floor:
                    end = i + snap + 1  # include the boundary char

            piece = text[i:end].strip()
            if piece:
                chunks.append(
                    Chunk(
                        index=idx,
                        content=piece,
                        metadata={
                            "page": page.page,
                            "section": _detect_section(piece, page_section_fallback),
                            "source_file": source_file,
                        },
                    )
                )
                idx += 1
            if end == len(text):
                break
            i += stride

    return chunks


# ────────────────────────── embedding ──────────────────────────

def _embed_chunks(texts: list[str]) -> list[Optional[list[float]]]:
    """Embed chunk texts, returning one vector per text (or ``None`` on
    failure so the chunk is still persisted without an embedding).

    Embedding is a best-effort enrichment: a provider outage degrades
    retrieval to the keyword fallback rather than failing the whole
    document.
    """
    if not texts:
        return []
    provider = get_embedding_provider()
    try:
        vectors = provider.embed(texts)
        return list(vectors)
    except Exception as e:  # noqa: BLE001 — degrade, don't fail the pipeline
        log.warning("embed_chunks failed (%s); storing chunks without vectors.", e)
        return [None] * len(texts)


# ────────────────────────── enrichment ──────────────────────────

_SENT_SPLIT_RE = re.compile(r"(?<=[.!?])\s+")
_WORD_RE = re.compile(r"[A-Za-z][A-Za-z'-]{2,}")
_STOPWORDS = {
    "the", "and", "for", "are", "but", "not", "you", "all", "any", "can", "her",
    "was", "one", "our", "out", "day", "get", "has", "him", "his", "how", "man",
    "new", "now", "old", "see", "two", "way", "who", "boy", "did", "its", "let",
    "put", "say", "she", "too", "use", "this", "that", "with", "from", "they",
    "have", "will", "your", "what", "when", "which", "their", "there", "would",
    "about", "into", "than", "then", "them", "these", "those", "been", "were",
    "also", "such", "each", "more", "most", "some", "other", "only", "over",
}


def compute_checksum(data: bytes) -> str:
    """SHA-256 hex digest of the raw bytes — drives duplicate detection."""
    import hashlib

    return hashlib.sha256(data).hexdigest()


def summarize_text(text: str, *, max_chars: int = 600) -> str:
    """Deterministic extractive summary: the leading sentences trimmed to
    ``max_chars``. Always available (no AI dependency) so processing never
    blocks on a provider outage."""
    clean = normalize(text)
    if not clean:
        return ""
    out: list[str] = []
    total = 0
    for sentence in _SENT_SPLIT_RE.split(clean.replace("\n", " ")):
        s = sentence.strip()
        if not s:
            continue
        if total + len(s) > max_chars and out:
            break
        out.append(s)
        total += len(s) + 1
    summary = " ".join(out).strip()
    if len(summary) > max_chars:
        summary = summary[:max_chars].rsplit(" ", 1)[0] + "…"
    return summary


def derive_questions(text: str, pages: list[ExtractedPage]) -> list[str]:
    """Generate a few suggested questions from section headings + top
    keywords. Deterministic, so it degrades gracefully without AI."""
    questions: list[str] = []
    seen: set[str] = set()

    # 1) Section headings → "What does <section> cover?"
    for p in pages:
        if p.section:
            label = p.section.strip().rstrip(":")[:80]
            key = label.lower()
            if label and key not in seen:
                seen.add(key)
                questions.append(f"What does the section on “{label}” cover?")
        if len(questions) >= 3:
            break

    # 2) Top keywords → generic prompts.
    if len(questions) < 3:
        counts: dict[str, int] = {}
        for w in _WORD_RE.findall(text.lower()):
            if w in _STOPWORDS:
                continue
            counts[w] = counts.get(w, 0) + 1
        top = [w for w, _ in sorted(counts.items(), key=lambda kv: kv[1], reverse=True)[:6]]
        for w in top:
            if len(questions) >= 3:
                break
            if w not in seen:
                seen.add(w)
                questions.append(f"What does this document say about {w}?")

    if not questions:
        questions = ["What is this document about?", "Summarize the key points."]
    return questions[:4]


# ────────────────────────── pipeline ──────────────────────────

async def process_document(document_id: uuid.UUID) -> None:
    """Run the full pipeline for one document, in its own session.

    Designed to be called from ``BackgroundTasks`` *after* the request
    that uploaded the document returns. Owns its own ``AsyncSession``
    so the original request session can close cleanly.

    Always-on guarantees:
      * Sets ``status='processing'`` + ``processing_started_at`` when work begins.
      * On success → ``status='processed'`` + ``processing_completed_at`` + chunk rows.
      * On any exception → ``status='failed'`` + ``processing_error`` populated.
    """
    if AsyncSessionLocal is None:
        init_engine()
    from app.database.session import AsyncSessionLocal as Maker

    async with Maker() as session:  # type: ignore[misc]
        await _run_pipeline(session, document_id)


async def _run_pipeline(session: AsyncSession, document_id: uuid.UUID) -> None:
    doc = await session.get(Document, document_id)
    if doc is None or doc.deleted_at is not None:
        log.warning("doc_process_skip missing_or_deleted document_id=%s", document_id)
        return

    # Status → processing
    doc.status = DocumentStatus.processing
    doc.processing_started_at = datetime.now(timezone.utc)
    doc.processing_error = None
    await session.commit()

    try:
        # Read bytes from storage
        if storage.is_local_key(doc.s3_key):
            data = storage.local_path(doc.s3_key).read_bytes()
        else:
            # S3 path — fetch with boto3 to a bytes buffer
            buf = io.BytesIO()
            storage._client().download_fileobj(  # type: ignore[attr-defined]
                Bucket=storage._bucket(),  # type: ignore[attr-defined]
                Key=doc.s3_key,
                Fileobj=buf,
            )
            data = buf.getvalue()

        pages = extract_text(doc.filename, doc.file_type, data)
        chunks = chunk_pages(pages, source_file=doc.filename)

        # R2 enrichment — always-on, deterministic so it can't fail the run.
        full_text = "\n\n".join(p.text for p in pages if p.text)
        try:
            doc.checksum = compute_checksum(data)
            doc.summary = summarize_text(full_text)
            doc.suggested_questions = derive_questions(full_text, pages)
            word_count = len(_WORD_RE.findall(full_text))
            doc.doc_metadata = {
                **(doc.doc_metadata or {}),
                "pages": len(pages),
                "word_count": word_count,
                "char_count": len(full_text),
                "chunk_count": len(chunks),
            }
        except Exception as e:  # noqa: BLE001 — enrichment is best-effort
            log.warning("doc_enrich_failed doc=%s err=%s", doc.id, e)

        # Phase 9: embed every chunk. Failure to embed must not lose the
        # chunks — we degrade to NULL embeddings (chunk still searchable
        # by the keyword fallback) and log it.
        embeddings = _embed_chunks([c.content for c in chunks])

        # Clear any prior chunks (reprocess path) then bulk-insert new ones.
        await session.execute(
            delete(DocumentChunk).where(DocumentChunk.document_id == doc.id)
        )
        for c, emb in zip(chunks, embeddings):
            session.add(
                DocumentChunk(
                    document_id=doc.id,
                    organization_id=doc.organization_id,
                    project_id=doc.project_id,
                    knowledge_base_id=doc.knowledge_base_id,
                    chunk_index=c.index,
                    content=c.content,
                    chunk_metadata=c.metadata,
                    embedding=emb,
                )
            )

        doc.status = DocumentStatus.processed
        doc.processing_completed_at = datetime.now(timezone.utc)
        await session.commit()

        embedded = sum(1 for e in embeddings if e is not None)
        log.info(
            "doc_process_ok doc=%s chunks=%d embedded=%d ms=%d",
            doc.id,
            len(chunks),
            embedded,
            int(
                (doc.processing_completed_at - doc.processing_started_at).total_seconds()
                * 1000
            ),
        )
    except Exception as e:  # noqa: BLE001 — we want to capture *any* failure
        log.exception("doc_process_failed doc=%s err=%s", doc.id, e)
        # Re-fetch after rollback to keep the row writable.
        await session.rollback()
        fresh = await session.get(Document, document_id)
        if fresh is not None:
            fresh.status = DocumentStatus.failed
            fresh.processing_completed_at = datetime.now(timezone.utc)
            fresh.processing_error = f"{type(e).__name__}: {e}"[:1000]
            await session.commit()


async def chunk_count_for(session: AsyncSession, document_id: uuid.UUID) -> int:
    from sqlalchemy import func

    n = await session.scalar(
        select(func.count(DocumentChunk.id)).where(
            DocumentChunk.document_id == document_id
        )
    )
    return int(n or 0)
